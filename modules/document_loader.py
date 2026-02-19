from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def load_text_from_file(file):
    file_name = file.name.lower()

    # ---------- PDF ----------
    if file_name.endswith(".pdf"):
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            if page.extract_text():
                text += page.extract_text() + "\n"
        return text

    # ---------- TXT ----------
    elif file_name.endswith(".txt"):
        return file.read().decode("utf-8", errors="ignore")

    # ---------- DOCX ----------
    elif file_name.endswith(".docx"):
        doc = Document(file)
        return "\n".join([p.text for p in doc.paragraphs])

    # ---------- PPTX ----------
    elif file_name.endswith(".pptx"):
        prs = Presentation(file)
        slides_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides_text.append(shape.text)
        return "\n".join(slides_text)

    else:
        raise ValueError("Unsupported file type")
