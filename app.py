import streamlit as st
import uuid

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

import chromadb
import ollama

# ---------------- Page Config ----------------
st.set_page_config(
    page_title="SelfLearnApp – Subject-Wise AI Study Assistant",
    page_icon="📘",
    layout="wide"
)

# ---------------- ChromaDB Setup ----------------
chroma_client = chromadb.Client(
    chromadb.config.Settings(
        persist_directory="./vectordb",
        anonymized_telemetry=False
    )
)

collection = chroma_client.get_or_create_collection(
    name="selflearn_subject_collection"
)

# ---------------- Subjects ----------------
subjects = [
    "Research Methodology and IPR",
    "Advanced Software Testing",
    "Large Language Models",
    "Cloud Architecture and Computing",
    "Mini Project",
    "Professional Communication",
    "AI-Driven Cybersecurity",
    "Agile Software Development"
]

# ---------------- Load Text ----------------
def load_text_from_file(file):
    name = file.name.lower()

    if name.endswith(".pdf"):
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            if page.extract_text():
                text += page.extract_text() + "\n"
        return text

    elif name.endswith(".txt"):
        return file.read().decode("utf-8", errors="ignore")

    elif name.endswith(".docx"):
        doc = Document(file)
        return "\n".join(p.text for p in doc.paragraphs)

    elif name.endswith(".pptx"):
        prs = Presentation(file)
        slides = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides.append(shape.text)
        return "\n".join(slides)

    return ""

# ---------------- Chunk Text (Balanced) ----------------
def chunk_text(text, chunk_size=450):
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunk = " ".join(words[i:i + chunk_size])
        if len(chunk.strip()) > 40:
            chunks.append(chunk)
    return chunks

# ---------------- Store Chunks (Safe) ----------------
def store_chunks(chunks, subject, source):
    if not chunks:
        return

    ids = [str(uuid.uuid4()) for _ in chunks]
    metadatas = [{"subject": subject, "source": source}] * len(chunks)

    collection.add(
        ids=ids,
        documents=chunks,
        metadatas=metadatas
    )

# ---------------- Sidebar ----------------
st.sidebar.title("📘 SelfLearnApp")
menu = st.sidebar.radio(
    "Navigation",
    ["🏠 Home", "📂 Professor Upload", "🔍 Ask Questions", "📊 Admin"]
)

# ---------------- Home ----------------
if menu == "🏠 Home":
    st.title("🎓 SelfLearnApp – Offline Intelligent Learning Assistant")
    st.markdown("""
    ✔ Subject-wise academic assistant  
    ✔ Faculty uploads (PDF / DOCX / PPTX)  
    ✔ Offline LLM (Ollama)  
    ✔ High accuracy, low hallucination  
    """)

# ---------------- Professor Upload ----------------
elif menu == "📂 Professor Upload":
    st.header("👩‍🏫 Professor Upload Portal")

    subject = st.selectbox("Select Subject", subjects)

    uploaded_files = st.file_uploader(
        "Upload PDF / TXT / DOCX / PPTX files",
        type=["pdf", "txt", "docx", "pptx"],
        accept_multiple_files=True
    )

    if uploaded_files:
        for file in uploaded_files:
            text = load_text_from_file(file)

            if not text or len(text.strip()) < 80:
                st.warning(f"⚠️ {file.name} has insufficient readable content")
                continue

            chunks = chunk_text(text)
            store_chunks(chunks, subject, file.name)

            st.success(f"✅ Stored {len(chunks)} chunks from {file.name}")

# ---------------- Ask Questions (ACCURATE + FAST) ----------------
elif menu == "🔍 Ask Questions":
    st.header("🔍 Ask a Question")

    subject = st.selectbox("Select Subject to Study Today", subjects)
    query = st.text_input("Enter your question")

    if query:
        with st.spinner("🧠 Analyzing academic content..."):
            results = collection.query(
                query_texts=[query],
                n_results=3,          # ✅ accuracy restored
                where={"subject": subject}
            )

            docs = results.get("documents", [[]])[0]

            if not docs:
                st.warning("⚠️ Not available in uploaded material.")
            else:
                context = "\n".join(docs)[:3000]  # ✅ balanced context

                response = ollama.chat(
                    model="mistral",  # ✅ accurate academic model
                    messages=[
                        {
                            "role": "system",
                            "content": (
                                "You are an academic assistant. "
                                "Answer strictly from the given context only. "
                                "Do not add outside knowledge. "
                                "If the answer is missing, say: "
                                "'Not available in uploaded material.'"
                            )
                        },
                        {
                            "role": "user",
                            "content": f"Context:\n{context}\n\nQuestion:\n{query}"
                        }
                    ],
                    options={
                        "temperature": 0,
                        "num_predict": 200
                    }
                )

                st.subheader("✅ Answer")
                st.write(response["message"]["content"])

# ---------------- Admin ----------------
elif menu == "📊 Admin":
    st.subheader("📊 Admin Dashboard")
    st.write(f"Total stored chunks: {collection.count()}")
